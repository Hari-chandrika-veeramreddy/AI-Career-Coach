import os
import pptx
from utils.speech_analyzer import transcribe_audio, analyze_audio_metrics

def parse_pptx(file_bytes_or_path):
    """
    Extract structured content from PowerPoint presentation (.pptx).
    Returns slide details, slide titles, text per slide, and summary stats.
    """
    stats = {
        "slide_count": 0,
        "total_words": 0,
        "avg_words_per_slide": 0.0,
        "empty_slides": 0,
        "text_heavy_slides": 0,
        "slide_titles": [],
        "all_text": "",
        "slides_detail": []
    }

    try:
        if isinstance(file_bytes_or_path, str) and os.path.exists(file_bytes_or_path):
            prs = pptx.Presentation(file_bytes_or_path)
        else:
            from io import BytesIO
            prs = pptx.Presentation(BytesIO(file_bytes_or_path))

        stats["slide_count"] = len(prs.slides)
        all_text_list = []

        for i, slide in enumerate(prs.slides, 1):
            slide_title = f"Slide {i}"
            slide_text = ""
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_str = shape.text.strip()
                    if hasattr(shape, "is_placeholder") and shape.is_placeholder and shape.placeholder_format.idx == 0:
                        slide_title = text_str
                    slide_text += text_str + " "

            slide_words = len(slide_text.split())
            if slide_words == 0:
                stats["empty_slides"] += 1
            elif slide_words > 80:
                stats["text_heavy_slides"] += 1

            stats["total_words"] += slide_words
            stats["slide_titles"].append(slide_title)
            stats["slides_detail"].append({
                "slide_number": i,
                "title": slide_title,
                "word_count": slide_words,
                "text": slide_text.strip()
            })
            all_text_list.append(slide_text.strip())

        stats["all_text"] = " ".join(all_text_list).strip()
        if stats["slide_count"] > 0:
            stats["avg_words_per_slide"] = round(stats["total_words"] / stats["slide_count"], 1)

    except Exception as e:
        print(f"PPTX parsing error: {e}")

    return stats

def evaluate_presentation(pptx_stats, topic="", audio_file=None):
    """
    Comprehensive presentation evaluation combining PPT structure analysis
    and spoken presentation audio analysis.
    """
    # 1. Content & Structure Score
    slide_count = pptx_stats.get("slide_count", 0)
    avg_words = pptx_stats.get("avg_words_per_slide", 0.0)
    text_heavy = pptx_stats.get("text_heavy_slides", 0)
    empty_slides = pptx_stats.get("empty_slides", 0)

    content_score = 100.0
    structure_score = 100.0

    if slide_count == 0:
        content_score = 0.0
        structure_score = 0.0
    elif slide_count < 3 or slide_count > 25:
        structure_score -= 20.0
    
    if text_heavy > 0:
        content_score -= min(30.0, text_heavy * 10.0)
    if empty_slides > 0:
        structure_score -= min(30.0, empty_slides * 15.0)
    if avg_words > 100.0:
        content_score -= 15.0

    content_score = round(max(20.0, min(100.0, content_score)), 1)
    structure_score = round(max(20.0, min(100.0, structure_score)), 1)

    # 2. Audio & Speech Analysis (if audio provided)
    audio_transcript = ""
    comm_metrics = {
        "duration": 0.0,
        "wpm": 0.0,
        "pause_ratio": 0.0,
        "communication_score": 75.0,
        "feedback": "Audio delivery was not provided."
    }

    if audio_file is not None:
        audio_transcript = transcribe_audio(audio_file)
        if hasattr(audio_file, "seek"):
            audio_file.seek(0)
        comm_metrics = analyze_audio_metrics(audio_file, audio_transcript)

    # 3. Spoken vs Slide Coverage Comparison
    ppt_text = pptx_stats.get("all_text", "").lower()
    spoken_text = audio_transcript.lower()
    
    coverage_score = 75.0
    if ppt_text and spoken_text:
        ppt_words = set(w for w in ppt_text.split() if len(w) > 3)
        spoken_words = set(w for w in spoken_text.split() if len(w) > 3)
        if ppt_words:
            overlap = ppt_words.intersection(spoken_words)
            coverage_ratio = len(overlap) / len(ppt_words)
            coverage_score = round(min(100.0, coverage_ratio * 120.0), 1)

    # Combined Overall Presentation Score
    comm_score = comm_metrics["communication_score"]
    overall_score = round((content_score * 0.35) + (structure_score * 0.25) + (comm_score * 0.25) + (coverage_score * 0.15), 1)

    # Strengths, Weaknesses, and Improvement Plan
    strengths = []
    weaknesses = []
    improvement_plan = []

    if slide_count >= 4:
        strengths.append(f"Good slide count ({slide_count} slides) suitable for presentation structure.")
    if avg_words <= 60.0 and slide_count > 0:
        strengths.append("Slide text density is balanced and clean, preventing visual overload.")

    if text_heavy > 0:
        weaknesses.append(f"{text_heavy} slide(s) are text-heavy with >80 words. Bullet points should be concise.")
        improvement_plan.append("Break down wall-of-text slides into 3-5 concise bullet points per slide.")

    if empty_slides > 0:
        weaknesses.append(f"{empty_slides} slide(s) contain no text content.")
        improvement_plan.append("Remove or populate blank slides before presenting.")

    if comm_metrics["wpm"] > 175:
        weaknesses.append(f"Fast speaking pace ({comm_metrics['wpm']} WPM).")
        improvement_plan.append("Practice speaking at 130-150 WPM to give your audience time to absorb complex points.")
    elif comm_metrics["wpm"] < 100 and comm_metrics["wpm"] > 0:
        weaknesses.append(f"Slower speaking pace ({comm_metrics['wpm']} WPM).")
        improvement_plan.append("Increase vocal energy to keep the audience engaged throughout.")

    if comm_metrics["pause_ratio"] > 35.0:
        weaknesses.append("Frequent long pauses detected during audio presentation.")
        improvement_plan.append("Rehearse slide transitions to ensure seamless verbal delivery.")

    if not strengths:
        strengths.append("Presentation outline submitted successfully.")
    if not weaknesses:
        weaknesses.append("No major slide formatting flaws detected.")
    if not improvement_plan:
        improvement_plan.append("Maintain current slide clarity and practice verbal delivery timing.")

    return {
        "score": overall_score,
        "content_score": content_score,
        "structure_score": structure_score,
        "communication_score": comm_score,
        "coverage_score": coverage_score,
        "speech_rate": comm_metrics.get("wpm", 0.0),
        "pause_score": 100.0 - comm_metrics.get("pause_ratio", 0.0),
        "confidence_score": comm_score,
        "transcript": audio_transcript,
        "strengths": strengths,
        "weaknesses": weaknesses,
        "improvement_plan": improvement_plan,
        "pptx_stats": pptx_stats,
        "comm_metrics": comm_metrics
    }
